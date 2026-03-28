"""PowerPoint → plain text."""
import io

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_text_from_pptx(file_stream) -> str:
    """Slurp all slide text from a file-like PPTX."""
    if Presentation is None:
        raise RuntimeError("python-pptx required for PPTX extraction. pip install python-pptx")

    prs = Presentation(file_stream)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n\n".join(parts)


def extract_text_from_pptx_bytes(data: bytes) -> str:
    return extract_text_from_pptx(io.BytesIO(data))
